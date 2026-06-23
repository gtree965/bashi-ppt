import io
import sys
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "backend"))

from llm.prompts import build_notes_messages
from llm.client import _extract_text_from_reasoning


def _outline():
    return {
        "title": "T",
        "slides": [
            {"page_number": 1, "title": "ТаЄжбў", "content_points": ["a", "b"], "slide_type": "title"},
            {"page_number": 2, "title": "тєЁт«╣", "content_points": ["x", "y", "z"], "slide_type": "content"},
            {"page_number": 3, "title": "у╗ЊУ»Г", "content_points": ["m", "n"], "slide_type": "conclusion"},
        ],
    }


class TestNotesPrompt(unittest.TestCase):
    def test_includes_duration_style_language(self):
        messages = build_notes_messages(
            _outline(),
            output_language="zh",
            duration_minutes=10,
            style="classroom",
        )
        system, user = messages[0]["content"], messages[1]["content"]
        self.assertIn("У»ЙтаѓУ«▓УДБ", system)   # style hint
        self.assertIn("10 тѕєжњЪ", user)       # duration budget
        self.assertIn("угг2жАх", user)         # per-slide listing

    def test_article_included_when_present(self):
        messages = build_notes_messages(
            _outline(),
            output_language="zh",
            duration_minutes=5,
            style="formal",
            article="ТќЄуФатєЁт«╣XYZ",
        )
        self.assertIn("ТќЄуФатєЁт«╣XYZ", messages[1]["content"])

    def test_grounded_notes_include_per_slide_fact_mapping(self):
        outline = _outline()
        outline["slides"][0]["fact_ids"] = [1]
        outline["slides"][1]["fact_ids"] = [2, 3]
        outline["slides"][2]["fact_ids"] = [3]
        messages = build_notes_messages(
            outline,
            output_language="zh",
            duration_minutes=10,
            style="classroom",
            mode="grounded",
            fact_table=[
                {"id": 1, "text": "С║Іт«ъСИђ"},
                {"id": 2, "text": "С║Іт«ъС║ї"},
                {"id": 3, "text": "С║Іт«ъСИЅ"},
            ],
        )
        system, user = messages[0]["content"], messages[1]["content"]
        self.assertIn("Т»ЈжАхтЈфУЃйт▒Ћт╝ђУ»ЦжАхТаЄТ│еуџёсђљС║Іт«ъу╝ќтЈисђЉ", system)
        self.assertIn("угг2жАх [content] тєЁт«╣№╝џx№╝Џy№╝Џz№╝ЏС║Іт«ъу╝ќтЈи№╝џ2сђЂ3", user)

    def test_reasoning_salvage_for_notes(self):
        text = 'think... {"notes":["n1","n2"]} end'
        self.assertEqual(_extract_text_from_reasoning(text, ('"notes"',)), '{"notes":["n1","n2"]}')
        # The default outline-key salvage must NOT grab the notes object.
        self.assertIsNone(_extract_text_from_reasoning(text))


class TestNotesParse(unittest.TestCase):
    def test_parse_variants(self):
        from app import _parse_notes
        self.assertEqual(_parse_notes('{"notes":["a","b"]}'), ["a", "b"])
        self.assertEqual(_parse_notes('["a","b"]'), ["a", "b"])
        self.assertFalse(_parse_notes("not json at all"))  # None or [] Рєњ both falsy


class TestNotesRendering(unittest.TestCase):
    def test_notes_embedded_per_slide(self):
        from pptx import Presentation
        from renderer.engine import PPTXRenderer

        outline = {
            "title": "T",
            "slides": [
                {"page_number": 1, "title": "ТаЄжбў", "content_points": ["a", "b"], "slide_type": "title", "notes": "т╝ђтю║уЎй"},
                {"page_number": 2, "title": "тєЁт«╣", "content_points": ["x", "y", "z"], "slide_type": "content", "notes": "У«▓У┐ЎжАх"},
                {"page_number": 3, "title": "у╗ЊУ»Г", "content_points": ["m", "n"], "slide_type": "conclusion"},
            ],
        }
        prs = Presentation(io.BytesIO(PPTXRenderer("teaching").render(outline)))
        self.assertEqual(prs.slides[0].notes_slide.notes_text_frame.text, "т╝ђтю║уЎй")
        self.assertEqual(prs.slides[1].notes_slide.notes_text_frame.text, "У«▓У┐ЎжАх")
        # A slide without notes should not get a notes slide created.
        self.assertFalse(prs.slides[2].has_notes_slide)


class TestNotesEndpoint(unittest.TestCase):
    def setUp(self):
        import app
        self.app_module = app
        self.client = app.app.test_client()

    def _body(self, **over):
        body = {
            "outline": _outline(),
            "output_language": "zh",
            "duration": 10,
            "style": "classroom",
        }
        body.update(over)
        return body

    def _fake(self, raw):
        from llm.client import LLMGenerationResult
        return LLMGenerationResult(raw_text=raw, elapsed_seconds=0.1, llm_model="test")

    def test_success(self):
        from unittest.mock import patch
        with patch.object(self.app_module, "generate_speaker_notes", return_value=self._fake('{"notes":["a","b","c"]}')):
            resp = self.client.post("/api/generate-notes", json=self._body())
        self.assertEqual(resp.status_code, 200)
        data = resp.get_json()
        self.assertEqual(data["notes"], ["a", "b", "c"])
        self.assertNotIn("warnings", data)

    def test_unparseable_returns_502(self):
        from unittest.mock import patch
        with patch.object(self.app_module, "generate_speaker_notes", return_value=self._fake("not json")):
            resp = self.client.post("/api/generate-notes", json=self._body())
        self.assertEqual(resp.status_code, 502)

    def test_grounded_notes_reject_incomplete_fact_mapping(self):
        from unittest.mock import patch

        outline = _outline()
        outline["slides"][0]["fact_ids"] = [1]
        outline["slides"][1]["fact_ids"] = []
        outline["slides"][2]["fact_ids"] = [1]
        with patch.object(self.app_module, "generate_speaker_notes") as generate:
            resp = self.client.post(
                "/api/generate-notes",
                json=self._body(
                    outline=outline,
                    generation_mode="grounded",
                    fact_table=[
                        {"id": 1, "text": "С║Іт«ъСИђ"},
                        {"id": 2, "text": "С║Іт«ъС║ї"},
                    ],
                ),
            )
        self.assertEqual(resp.status_code, 422)
        audit = resp.get_json()["grounding_audit"]
        self.assertEqual(audit["missing_fact_ids"], [2])
        self.assertEqual(audit["ungrounded_content_pages"], [2])
        generate.assert_not_called()

    def test_length_mismatch_warns(self):
        from unittest.mock import patch
        with patch.object(self.app_module, "generate_speaker_notes", return_value=self._fake('{"notes":["only one"]}')):
            resp = self.client.post("/api/generate-notes", json=self._body())
        self.assertEqual(resp.status_code, 200)
        data = resp.get_json()
        self.assertEqual(len(data["notes"]), 3)
        self.assertTrue(data.get("warnings"))

    def test_invalid_duration_returns_422(self):
        resp = self.client.post("/api/generate-notes", json=self._body(duration=7))
        self.assertEqual(resp.status_code, 422)


if __name__ == "__main__":
    unittest.main()
