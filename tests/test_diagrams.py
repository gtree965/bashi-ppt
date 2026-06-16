import base64
import io
import sys
import unittest
from pathlib import Path

# Add backend to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "backend"))

from llm.outline_parser import parse_outline
from renderer.engine import PPTXRenderer


def _make_png_data_url() -> str:
    """Create a tiny valid PNG and return it as a base64 data-URL."""
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (320, 240), color=(255, 255, 255)).save(buffer, format="PNG")
    encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


class TestDiagramParsing(unittest.TestCase):
    def test_diagram_preserved_on_content_slide(self):
        raw = """
        {
          "title": "搜索引擎如何工作",
          "slides": [
            {"page_number": 1, "title": "搜索引擎", "content_points": ["要点一", "要点二"], "slide_type": "title"},
            {"page_number": 2, "title": "工作流程", "content_points": ["抓取", "索引", "排序"],
             "slide_type": "content", "diagram": "flowchart LR; A-->B-->C", "chart_config": {"x": 1}},
            {"page_number": 3, "title": "排序原理", "content_points": ["相关性", "权重", "反馈"], "slide_type": "content"},
            {"page_number": 4, "title": "总结", "content_points": ["快速", "相关"], "slide_type": "conclusion"}
          ]
        }
        """
        result = parse_outline(raw)
        slides = result.outline["slides"]

        # diagram preserved on the content slide
        self.assertEqual(slides[1].get("diagram"), "flowchart LR; A-->B-->C")
        # forbidden chart_config still stripped
        self.assertNotIn("chart_config", slides[1])

    def test_diagram_dropped_on_title_and_conclusion(self):
        raw = """
        {
          "title": "演示",
          "slides": [
            {"page_number": 1, "title": "标题", "content_points": ["a", "b"],
             "slide_type": "title", "diagram": "flowchart LR; X-->Y"},
            {"page_number": 2, "title": "内容", "content_points": ["a", "b", "c"], "slide_type": "content"},
            {"page_number": 3, "title": "内容二", "content_points": ["a", "b", "c"], "slide_type": "content"},
            {"page_number": 4, "title": "结语", "content_points": ["a", "b"],
             "slide_type": "conclusion", "diagram": "flowchart LR; X-->Y"}
          ]
        }
        """
        result = parse_outline(raw)
        slides = result.outline["slides"]
        self.assertIsNone(slides[0].get("diagram"))
        self.assertIsNone(slides[-1].get("diagram"))


class TestDiagramRendering(unittest.TestCase):
    def _render(self, outline):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        pptx_bytes = PPTXRenderer("teaching").render(outline)
        prs = Presentation(io.BytesIO(pptx_bytes))
        return prs, MSO_SHAPE_TYPE

    def _base_outline(self, content_slide):
        return {
            "title": "Diagram Test",
            "slides": [
                {"page_number": 1, "title": "Title", "content_points": ["a", "b"], "slide_type": "title"},
                content_slide,
                {"page_number": 3, "title": "End", "content_points": ["a", "b"], "slide_type": "conclusion"},
            ],
        }

    def test_diagram_image_embedded(self):
        outline = self._base_outline({
            "page_number": 2, "title": "Flow", "content_points": ["x", "y", "z"],
            "slide_type": "content", "diagram_image": _make_png_data_url(),
        })
        prs, MSO_SHAPE_TYPE = self._render(outline)
        content_slide = prs.slides[1]
        pictures = [s for s in content_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        self.assertEqual(len(pictures), 1, "diagram_image should be embedded as one picture")

    def test_image_url_takes_precedence_over_diagram(self):
        # When both are set, no diagram should be decoded/embedded if image_url wins.
        # Use an unreachable image_url so download fails; the diagram must NOT fill in,
        # confirming precedence is decided before any diagram decode.
        outline = self._base_outline({
            "page_number": 2, "title": "Flow", "content_points": ["x", "y", "z"],
            "slide_type": "content",
            "image_url": "http://127.0.0.1:9/nonexistent.png",
            "diagram_image": _make_png_data_url(),
        })
        prs, MSO_SHAPE_TYPE = self._render(outline)
        content_slide = prs.slides[1]
        pictures = [s for s in content_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        self.assertEqual(len(pictures), 0, "diagram must not be used when image_url is present")


if __name__ == "__main__":
    unittest.main()
