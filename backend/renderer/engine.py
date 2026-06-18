import json
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import logging

from .theme import THEMES
from .slide_layouts import LAYOUTS, SLIDE_WIDTH, SLIDE_HEIGHT
from .utils import hex_to_rgb, set_font
import config
from config import TEMPLATES_DIR

logger = logging.getLogger("slideforge")


def set_chinese_font(run, font_name="Microsoft YaHei", size_pt=20, bold=False, color: RGBColor = None):
    """Backward-compatible wrapper around set_font for Chinese text."""
    set_font(run, font_name=font_name, size_pt=size_pt, bold=bold, color=color, script="zh")

def estimate_text_height(content_points, font_size_pt, box_width_inches, line_spacing=1.5):
    total_lines = 0
    # Approximate ratio for Chinese chars vs pts.
    chars_per_line = int(box_width_inches * 72 / font_size_pt) 

    for point in content_points:
        lines_needed = max(1, -(-len(point) // chars_per_line)) if chars_per_line > 0 else 1
        total_lines += lines_needed

    height_inches = total_lines * font_size_pt * line_spacing / 72
    return height_inches

def _fit_text_in_box(max_height_inches, box_width_inches, content_points, base_font_size=20, min_font_size=14):
    font_size = base_font_size
    line_spacing = 1.5
    points = content_points.copy()
    layout_warning = False

    while font_size >= min_font_size:
        if estimate_text_height(points, font_size, box_width_inches, line_spacing) <= max_height_inches:
            return font_size, line_spacing, points, layout_warning
        font_size -= 2

    # Still overflowing at min size, try reducing spacing
    line_spacing = 1.2
    if estimate_text_height(points, font_size, box_width_inches, line_spacing) <= max_height_inches:
        return font_size, line_spacing, points, layout_warning

    # Still overflowing, truncate text
    layout_warning = True
    max_iterations = 1000
    iterations = 0
    while points and estimate_text_height(points, font_size, box_width_inches, line_spacing) > max_height_inches:
        iterations += 1
        if iterations > max_iterations:
            logger.warning("Reached maximum iterations in _fit_text_in_box. Force truncating.")
            points.pop()
            continue
            
        if len(points[-1]) > 5:
            points[-1] = points[-1][:-5] + "..."
        else:
            points.pop()
            
    return font_size, line_spacing, points, layout_warning

class PPTXRenderer:
    def __init__(self, template_id: str, theme_id: str | None = None):
        self.template_id = template_id
        template_path = TEMPLATES_DIR / f"{template_id}.json"
        
        if not template_path.exists():
            template_path = TEMPLATES_DIR / "teaching.json"
            
        with open(template_path, 'r', encoding='utf-8') as f:
            self.template_config = json.load(f)
            
        if not theme_id:
            theme_id = self.template_config.get("theme", "clean_blue")
            
        self.theme = THEMES.get(theme_id, THEMES["clean_blue"])
        self.layout_mapping = self.template_config.get("layout_mapping", {})
        
        self.colors = self.theme["colors"]
        self.fonts = self.theme["fonts"]
        self.decorations = self.theme["decorations"]

    def render(self, outline: dict, bullet_style: str = "dot") -> bytes:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        
        slides_data = outline.get("slides", [])
        total_slides = len(slides_data)
        
        for idx, slide_data in enumerate(slides_data):
            self._render_slide(prs, slide_data, idx + 1, total_slides, bullet_style=bullet_style)
            
        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
        
    def _apply_speaker_notes(self, slide, slide_data: dict):
        """Write speaker notes into the slide's PowerPoint notes pane, if present."""
        notes = (slide_data.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _create_blank_slide(self, prs):
        # Index 6 is typically 'Blank'
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        bg_color = hex_to_rgb(self.colors["background"])
        fill.fore_color.rgb = bg_color
        return slide

    def _render_slide(self, prs, slide_data: dict, current: int, total: int, bullet_style: str = "dot"):
        slide_type = slide_data.get("slide_type", "content")
        layout_id = self.layout_mapping.get(slide_type)
        if slide_type == "title":
            layout_id = layout_id or "TitleCenterLayout"
            self._render_title_slide(prs, slide_data, layout_id, current, total)
        elif slide_type == "conclusion":
            layout_id = layout_id or "ConclusionLayout"
            self._render_conclusion_slide(prs, slide_data, layout_id, current, total)
        else:
            layout_id = layout_id or "ContentBulletLayout"
            self._render_content_slide(prs, slide_data, layout_id, current, total, bullet_style=bullet_style)
            
    def _render_title_slide(self, prs, slide_data: dict, layout_id: str, current: int, total: int):
        slide = self._create_blank_slide(prs)
        self._apply_speaker_notes(slide, slide_data)
        layout = LAYOUTS.get(layout_id, LAYOUTS["TitleCenterLayout"])
        
        # Title
        t_box = layout["title_box"]
        txBox = slide.shapes.add_textbox(t_box["left"], t_box["top"], t_box["width"], t_box["height"])
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data.get("title", "")
        set_chinese_font(run, font_name=self.fonts["title"], size_pt=self.fonts.get("title_size", 44), 
                         bold=True, color=hex_to_rgb(self.colors["primary"]))
                         
        # Subtitle (Points)
        points = slide_data.get("content_points", [])
        if points:
            s_box = layout["subtitle_box"]
            sxBox = slide.shapes.add_textbox(s_box["left"], s_box["top"], s_box["width"], s_box["height"])
            sf = sxBox.text_frame
            sf.word_wrap = True
            p = sf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "\n".join(points)
            set_chinese_font(run, font_name=self.fonts["body"], size_pt=self.fonts.get("body_size", 20), 
                             color=hex_to_rgb(self.colors["secondary"]))

        # Decorative line
        if self.decorations.get("title_underline", False) and "accent_line" in layout:
            line_box = layout["accent_line"]
            shape = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                line_box["left"], line_box["top"], line_box["width"], line_box["height"]
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(self.colors["accent"])
            shape.line.fill.background()
            
    def _download_image(self, url: str) -> str | None:
        """Download an image to a temporary file, returning the path."""
        import urllib.request
        import tempfile
        import ssl
        from pathlib import Path
        
        try:
            suffix = Path(url).suffix or ".jpg"
            if "?" in suffix:
                suffix = suffix.split("?")[0]
            if suffix.lower() not in [".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"]:
                suffix = ".jpg"
                
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            
            # Disable SSL verification for proxies/VPNs if requested
            if not getattr(config, "VERIFY_SSL", True):
                ctx = ssl.create_default_context()
                ctx.check_hostname = False
                ctx.verify_mode = ssl.CERT_NONE
            else:
                ctx = None
            
            headers = {'User-Agent': 'Mozilla/5.0'}
            req = urllib.request.Request(url, headers=headers)
            
            with urllib.request.urlopen(req, context=ctx, timeout=10) as response, open(temp_file.name, 'wb') as out_file:
                out_file.write(response.read())
                
            return temp_file.name
        except Exception as e:
            logger.error(f"Failed to download image {url}: {e}")
            return None

    def _decode_data_url(self, data_url: str) -> str | None:
        """Decode a base64 image data-URL to a temporary PNG file, returning the path.

        Used for hand-drawn diagrams rendered in the browser and sent as
        ``data:image/png;base64,...``. ``_download_image`` cannot handle these
        because it speaks urllib/HTTP, not data URLs.
        """
        import base64
        import tempfile

        try:
            payload = data_url
            if payload.startswith("data:"):
                # Strip the "data:image/png;base64," prefix.
                _, _, payload = payload.partition(",")
            if not payload:
                return None

            image_bytes = base64.b64decode(payload)
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            try:
                temp_file.write(image_bytes)
            finally:
                temp_file.close()
            return temp_file.name
        except Exception as e:
            logger.error(f"Failed to decode diagram data-URL: {e}")
            return None

    def _render_content_slide(self, prs, slide_data: dict, layout_id: str, current: int, total: int, bullet_style: str = "dot"):
        slide = self._create_blank_slide(prs)
        self._apply_speaker_notes(slide, slide_data)
        layout = LAYOUTS.get(layout_id, LAYOUTS["ContentBulletLayout"])
        
        # Title
        t_box = layout["title_box"]
        txBox = slide.shapes.add_textbox(t_box["left"], t_box["top"], t_box["width"], t_box["height"])
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data.get("title", "")
        set_chinese_font(run, font_name=self.fonts["title"], size_pt=self.fonts.get("title_size", 32), 
                         bold=True, color=hex_to_rgb(self.colors["primary"]))
                         
        # Points
        points = slide_data.get("content_points", [])
        bullet_prefixes = {
            "dot": "•  ",
            "square": "▪  ",
            "arrow": "➢  ",
            "dash": "-  ",
            "none": "",
        }
        bullet_prefix = bullet_prefixes.get(bullet_style, "•  ")
        bulleted_points = [f"{bullet_prefix}{pt}" if pt else pt for pt in points]
        
        # Check if slide has a graphic. A search image takes precedence over a
        # rendered diagram, since the right column only holds one graphic.
        image_url = slide_data.get("image_url")
        diagram_image = slide_data.get("diagram_image")
        temp_img_path = None
        has_image = False

        if image_url:
            temp_img_path = self._download_image(image_url)
        elif diagram_image:
            temp_img_path = self._decode_data_url(diagram_image)
        if temp_img_path:
            has_image = True
                
        if has_image:
            # Side-by-side coordinates
            c_box_left = Inches(1.0)
            c_box_top = Inches(1.8)
            c_box_w = Inches(6.2)
            c_box_h = Inches(4.8)
            
            image_left = Inches(7.8)
            image_top = Inches(1.8)
            image_w = Inches(4.5)
            image_h = Inches(4.5)
        else:
            c_box = layout["content_box"]
            c_box_left = c_box["left"]
            c_box_top = c_box["top"]
            c_box_w = c_box["width"]
            c_box_h = c_box["height"]
        
        # Calculate overflow
        box_w = c_box_w.inches
        box_h = c_box_h.inches
        f_size, line_spacing, fitted_points, layout_warning = _fit_text_in_box(
            box_h, box_w, bulleted_points, base_font_size=self.fonts.get("body_size", 20)
        )
        if layout_warning:
            logger.warning(f"Text overflow on slide {current}, truncated content.")
            
        if fitted_points:
            cxBox = slide.shapes.add_textbox(c_box_left, c_box_top, c_box_w, c_box_h)
            cf = cxBox.text_frame
            cf.word_wrap = True
            
            for i, point in enumerate(fitted_points):
                p = cf.paragraphs[0] if i == 0 else cf.add_paragraph()
                p.text = point
                p.level = 0
                p.line_spacing = line_spacing
                p.space_after = Pt(8)
                run = p.runs[0] if p.runs else p.add_run()
                if not p.runs:
                    run.text = point
                set_chinese_font(run, font_name=self.fonts["body"], size_pt=f_size, 
                                 color=hex_to_rgb(self.colors["text"]))

        # Add image shape if present
        if has_image and temp_img_path:
            import os
            from PIL import Image
            try:
                # Open image using Pillow to check dimensions and aspect ratio
                with Image.open(temp_img_path) as img:
                    orig_w, orig_h = img.size
                
                # Default box boundaries: 4.5" x 4.5" centered in the right column
                max_w = Inches(4.5)
                max_h = Inches(4.5)
                box_left = Inches(7.8)
                box_top = Inches(1.8)
                
                # Keep aspect ratio and center inside the bounding box
                ratio = orig_w / orig_h
                if ratio > 1.0:  # Landscape
                    final_w = max_w
                    final_h = Inches(4.5 / ratio)
                else:  # Portrait / Square
                    final_h = max_h
                    final_w = Inches(4.5 * ratio)
                
                final_left = box_left + (max_w - final_w) / 2
                final_top = box_top + (max_h - final_h) / 2
                
                slide.shapes.add_picture(temp_img_path, final_left, final_top, width=final_w, height=final_h)
            except Exception as e:
                logger.error(f"Failed to add picture to slide {current}: {e}")
                # Fallback to default stretched placement if Pillow analysis fails
                try:
                    slide.shapes.add_picture(temp_img_path, image_left, image_top, width=image_w, height=image_h)
                except Exception as fallback_err:
                    logger.error(f"Fallback add picture also failed: {fallback_err}")
            finally:
                try:
                    os.unlink(temp_img_path)
                except Exception:
                    pass

        # Optional accent bar (only if no image, to avoid clutter)
        if not has_image and self.decorations.get("accent_bar", False) and "accent_bar" in layout:
            bar_box = layout["accent_bar"]
            shape = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                bar_box["left"], bar_box["top"], bar_box["width"], bar_box["height"]
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(self.colors["primary"])
            shape.line.fill.background()

        # Page Number
        self._add_slide_number(slide, layout, current, total)
        
    def _render_conclusion_slide(self, prs, slide_data: dict, layout_id: str, current: int, total: int):
        slide = self._create_blank_slide(prs)
        self._apply_speaker_notes(slide, slide_data)
        layout = LAYOUTS.get(layout_id, LAYOUTS["ConclusionLayout"])
        
        # Title
        t_box = layout["title_box"]
        txBox = slide.shapes.add_textbox(t_box["left"], t_box["top"], t_box["width"], t_box["height"])
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data.get("title", "")
        set_chinese_font(run, font_name=self.fonts["title"], size_pt=self.fonts.get("title_size", 44), 
                         bold=True, color=hex_to_rgb(self.colors["primary"]))

        # Points
        points = slide_data.get("content_points", [])
        if points:
            p_box = layout["points_box"]
            pxBox = slide.shapes.add_textbox(p_box["left"], p_box["top"], p_box["width"], p_box["height"])
            pf = pxBox.text_frame
            pf.word_wrap = True
            for i, point in enumerate(points):
                p = pf.paragraphs[0] if i == 0 else pf.add_paragraph()
                p.text = point
                p.level = 0
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                if not p.runs:
                    run.text = point
                set_chinese_font(run, font_name=self.fonts["body"], size_pt=self.fonts.get("body_size", 20), 
                                 color=hex_to_rgb(self.colors["secondary"]))

        # Page Number
        self._add_slide_number(slide, layout, current, total)

    def _add_slide_number(self, slide, layout, current: int, total: int):
        if "page_number" not in layout:
            return
            
        fmt = self.decorations.get("slide_number_format", "{current}")
        text = fmt.format(current=current, total=total)
        
        num_box = layout["page_number"]
        bxBox = slide.shapes.add_textbox(num_box["left"], num_box["top"], num_box["width"], num_box["height"])
        tf = bxBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        set_chinese_font(run, font_name=self.fonts.get("body", "Arial"), size_pt=self.fonts.get("page_number_size", 10), 
                         color=hex_to_rgb(self.colors["page_number"]))
