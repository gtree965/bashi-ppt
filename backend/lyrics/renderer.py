"""
Lyrics PPTX Renderer — optimized for worship projection.

Key differences from the presentation renderer:
  - No title bars, no bullet points, no page numbers
  - Large centered text on dark backgrounds
  - Chorus highlighting in accent color
  - Optional bilingual rendering (primary + secondary lines)
"""

from io import BytesIO

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from lyrics.chinese_script import convert_text
from renderer.utils import hex_to_rgb, set_font
from renderer.slide_layouts import SLIDE_WIDTH, SLIDE_HEIGHT

# =====================================================================
# Theme & font configuration
# =====================================================================

LYRICS_THEMES = {
    "classic_dark": {
        "name": "经典黑底白字",
        "background": "000000",
        "text_color": "FFFFFF",
        "chorus_color": "FFD700",
        "secondary_text_color": "B0BEC5",
    },
    "deep_blue": {
        "name": "深蓝渐变",
        "background": "0C1445",
        "text_color": "FFFFFF",
        "chorus_color": "81D4FA",
        "secondary_text_color": "90CAF9",
    },
    "warm_dark": {
        "name": "暖色深底",
        "background": "1A0A00",
        "text_color": "FFF8E1",
        "chorus_color": "FFB74D",
        "secondary_text_color": "D7CCC8",
    },
}

LANGUAGE_FONT_SIZES = {
    "single":              {"zh": 40, "ko": 40, "ja": 40, "latin": 38},
    "bilingual_primary":   {"zh": 36, "ko": 36, "ja": 36, "latin": 34},
    "bilingual_secondary": {"zh": 28, "ko": 28, "ja": 28, "latin": 26},
}

LANGUAGE_FONTS = {
    "zh": "Microsoft YaHei",
    "ko": "Malgun Gothic",
    "ja": "Yu Gothic",
    "latin": "Arial",
}

# Text box position (centered with margin)
_TEXT_LEFT = 1.0     # inches
_TEXT_TOP = 1.0      # inches
_TEXT_WIDTH = 11.333  # inches
_TEXT_HEIGHT = 5.5    # inches

_TITLE_SIZE = 54
_AMEN_SIZE = 60


def _get_script_for_lang(lang_code: str) -> str:
    """Map a language code to a script key for font/size lookup."""
    _LANG_TO_SCRIPT = {
        "zh": "zh", "ko": "ko", "ja": "ja",
        "en": "latin", "fr": "latin", "es": "latin",
        "id": "latin", "pt": "latin", "de": "latin",
    }
    return _LANG_TO_SCRIPT.get(lang_code, "latin")


def _set_vertical_center(text_frame):
    """Set vertical anchor to middle on a text frame via XML."""
    txBody = text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = txBody.makeelement(qn('a:bodyPr'), {})
        txBody.insert(0, bodyPr)
    bodyPr.set('anchor', 'ctr')


class LyricsPPTXRenderer:
    """Generates worship-optimized lyrics PPTX files."""

    def render(
        self,
        slides_data: list[dict],
        title: str = "",
        theme: str = "classic_dark",
        language_mode: str = "single",
        language_config: dict | None = None,
        add_title_slide: bool = True,
        add_amen_slide: bool = True,
        bilingual_pairs: list[list[tuple[str, str]]] | None = None,
        font_family: str | None = None,
        font_size_adjustment: int = 0,
        line_spacing: float = 1.5,
    ) -> bytes:
        """
        Render lyrics slides to PPTX bytes.

        Parameters
        ----------
        slides_data : list of {"lines": [...], "is_chorus": bool, "type": "lyrics"}
        title : song title
        theme : theme key
        language_mode : "single" or "bilingual"
        language_config : {"primary": "zh", "secondary": "en", ...}
        add_title_slide : whether to add a title slide at the start
        add_amen_slide : whether to add an "Amen" slide at the end
        bilingual_pairs : for bilingual mode, list of [(primary, secondary), ...] per slide
        """
        theme_cfg = LYRICS_THEMES.get(theme, LYRICS_THEMES["classic_dark"])
        lang_cfg = language_config or {"primary": "zh"}

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        if add_title_slide:
            self._render_title_slide(prs, title, theme_cfg, lang_cfg, font_family, font_size_adjustment)

        if language_mode == "bilingual" and bilingual_pairs:
            for i, slide_data in enumerate(slides_data):
                pairs = bilingual_pairs[i] if i < len(bilingual_pairs) else []
                self._render_bilingual_slide(prs, pairs, theme_cfg, slide_data.get("is_chorus", False), lang_cfg, font_family, font_size_adjustment, line_spacing)
        else:
            for slide_data in slides_data:
                self._render_lyrics_slide(
                    prs, slide_data["lines"], theme_cfg,
                    slide_data.get("is_chorus", False), lang_cfg, font_family, font_size_adjustment, line_spacing
                )

        if add_amen_slide:
            self._render_amen_slide(prs, theme_cfg, lang_cfg, font_family, font_size_adjustment)

        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _create_blank_slide(self, prs, theme_cfg: dict):
        """Create a blank slide with dark background."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(theme_cfg["background"])
        return slide

    def _render_title_slide(self, prs, title: str, theme_cfg: dict, lang_cfg: dict, font_family: str | None = None, font_size_adjustment: int = 0):
        """Title slide: song name centered large."""
        from pptx.util import Inches

        slide = self._create_blank_slide(prs, theme_cfg)
        txBox = slide.shapes.add_textbox(
            Inches(_TEXT_LEFT), Inches(2.5),
            Inches(_TEXT_WIDTH), Inches(2.5),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_vertical_center(tf)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title

        primary_script = _get_script_for_lang(lang_cfg.get("primary", "zh"))
        font_name = font_family or LANGUAGE_FONTS.get(primary_script, "Arial")
        set_font(run, font_name=font_name, size_pt=_TITLE_SIZE + font_size_adjustment, bold=True,
                 color=hex_to_rgb(theme_cfg["text_color"]), script=primary_script)

    def _render_lyrics_slide(
        self, prs, lines: list[str], theme_cfg: dict,
        is_chorus: bool, lang_cfg: dict, font_family: str | None = None, font_size_adjustment: int = 0, line_spacing: float = 1.5
    ):
        """Single-language lyrics slide: centered large text."""
        from pptx.util import Inches

        slide = self._create_blank_slide(prs, theme_cfg)
        txBox = slide.shapes.add_textbox(
            Inches(_TEXT_LEFT), Inches(_TEXT_TOP),
            Inches(_TEXT_WIDTH), Inches(_TEXT_HEIGHT),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_vertical_center(tf)

        primary_script = _get_script_for_lang(lang_cfg.get("primary", "zh"))
        font_name = font_family or LANGUAGE_FONTS.get(primary_script, "Arial")
        font_size = LANGUAGE_FONT_SIZES["single"].get(primary_script, 38) + font_size_adjustment
        text_color = theme_cfg["chorus_color"] if is_chorus else theme_cfg["text_color"]

        for i, line_text in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(4)
            p.line_spacing = line_spacing
            run = p.add_run()
            run.text = line_text
            set_font(run, font_name=font_name, size_pt=font_size, bold=True,
                     color=hex_to_rgb(text_color), script=primary_script)

    def _render_bilingual_slide(
        self, prs, pairs: list[tuple[str, str]], theme_cfg: dict,
        is_chorus: bool, lang_cfg: dict, font_family: str | None = None, font_size_adjustment: int = 0, line_spacing: float = 1.5
    ):
        """Bilingual lyrics slide: primary (large) + secondary (smaller, gray) alternating."""
        from pptx.util import Inches

        slide = self._create_blank_slide(prs, theme_cfg)
        txBox = slide.shapes.add_textbox(
            Inches(_TEXT_LEFT), Inches(_TEXT_TOP),
            Inches(_TEXT_WIDTH), Inches(_TEXT_HEIGHT),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_vertical_center(tf)

        primary_script = _get_script_for_lang(lang_cfg.get("primary", "zh"))
        secondary_script = _get_script_for_lang(lang_cfg.get("secondary", "en"))

        primary_font = font_family or LANGUAGE_FONTS.get(primary_script, "Arial")
        secondary_font = font_family or LANGUAGE_FONTS.get(secondary_script, "Arial")

        primary_size = LANGUAGE_FONT_SIZES["bilingual_primary"].get(primary_script, 34) + font_size_adjustment
        secondary_size = LANGUAGE_FONT_SIZES["bilingual_secondary"].get(secondary_script, 26) + font_size_adjustment

        primary_color = theme_cfg["chorus_color"] if is_chorus else theme_cfg["text_color"]
        secondary_color = theme_cfg["secondary_text_color"]

        first_para = True
        for primary_text, secondary_text in pairs:
            # Primary line
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)
            p.line_spacing = line_spacing
            run = p.add_run()
            run.text = primary_text
            set_font(run, font_name=primary_font, size_pt=primary_size, bold=True,
                     color=hex_to_rgb(primary_color), script=primary_script)

            # Secondary line
            if secondary_text:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_after = Pt(8)
                p2.line_spacing = line_spacing
                run2 = p2.add_run()
                run2.text = secondary_text
                set_font(run2, font_name=secondary_font, size_pt=secondary_size, bold=False,
                         color=hex_to_rgb(secondary_color), script=secondary_script)

    def _render_amen_slide(self, prs, theme_cfg: dict, lang_cfg: dict, font_family: str | None = None, font_size_adjustment: int = 0):
        """Amen slide: large centered text."""
        from pptx.util import Inches

        slide = self._create_blank_slide(prs, theme_cfg)
        txBox = slide.shapes.add_textbox(
            Inches(_TEXT_LEFT), Inches(2.5),
            Inches(_TEXT_WIDTH), Inches(2.5),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_vertical_center(tf)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()

        primary_script = _get_script_for_lang(lang_cfg.get("primary", "zh"))
        if primary_script == "zh":
            amen_text = "阿们"
            run.text = convert_text(amen_text, lang_cfg.get("script_conversion", "original"))
        else:
            run.text = "Amen"

        font_name = font_family or LANGUAGE_FONTS.get(primary_script, "Arial")
        set_font(run, font_name=font_name, size_pt=_AMEN_SIZE + font_size_adjustment, bold=True,
                 color=hex_to_rgb(theme_cfg["text_color"]), script=primary_script)
